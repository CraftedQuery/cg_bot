"""
rag_chatbot.py – multi‑tenant RAG bot (v7.0, enhanced)
• Added GUI for command line interface
• Completed the frontend widget with advanced features
• Implemented authentication mechanism
• Added support for additional document formats
• Added alternative LLM providers
• Implemented real-time analytics
"""

from __future__ import annotations
import os
import sys
import json
import tempfile
import contextlib
import xml.etree.ElementTree as ET
import sqlite3
import datetime
import time
from pathlib import Path
from typing import List, Optional, Dict, Union, Any

# FastAPI imports
import requests
from bs4 import BeautifulSoup
from fastapi import FastAPI, HTTPException, Request, Query, Depends, status, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import PlainTextResponse, JSONResponse, HTMLResponse
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

# Authentication imports
from jose import JWTError, jwt
from passlib.context import CryptContext
from datetime import datetime, timedelta

# RAG-related imports
from langchain_openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS

# Google Drive imports
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# GUI imports for CLI
import typer
from rich.console import Console
from rich.panel import Panel
from rich.layout import Layout
from rich.table import Table
from rich.progress import Progress
from rich import box
import pandas as pd

# ───────────────────────── Basics ───────────────────────────
BASE_CONFIG_DIR = Path("configs")
BASE_CONFIG_DIR.mkdir(exist_ok=True)
BASE_STORE_DIR = Path("vector_store")
DEFAULT_TENANT = "public"
DEFAULT_AGENT = "default"
TEXT_SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=80)

# ──────────────────────── Logging DB ────────────────────────
DB_PATH = Path("chat_logs.db")
with sqlite3.connect(DB_PATH) as con:
    con.execute("""CREATE TABLE IF NOT EXISTS chat_logs(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        ts TEXT, tenant TEXT, agent TEXT, session_id TEXT,
        question TEXT, answer TEXT, sources TEXT,
        latency REAL, tokens_in INTEGER, tokens_out INTEGER,
        user_feedback INTEGER, user_ip TEXT)""")

# ────────────────────────── Auth ───────────────────────────
# Authentication configuration
SECRET_KEY = os.getenv("JWT_SECRET_KEY", "dev_secret_key_change_in_production")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30

# Password hashing
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

# User model
class User(BaseModel):
    username: str
    tenant: str
    role: str = "user"
    disabled: bool = False

class UserCreate(User):
    password: str

# Token models
class Token(BaseModel):
    access_token: str
    token_type: str
    
class TokenData(BaseModel):
    username: Optional[str] = None
    tenant: Optional[str] = None
    role: Optional[str] = None

# User database - this would be replaced with a proper DB in production
def get_users_db():
    users_file = Path("users.json")
    if not users_file.exists():
        # Create default admin user if no users file exists
        default_admin = {
            "admin": {
                "username": "admin",
                "tenant": "*",  # Wildcard for all tenants
                "role": "admin",
                "hashed_password": pwd_context.hash("admin"),
                "disabled": False
            }
        }
        users_file.write_text(json.dumps(default_admin, indent=2))
        return default_admin
    return json.loads(users_file.read_text())

def save_users_db(users_data):
    users_file = Path("users.json")
    users_file.write_text(json.dumps(users_data, indent=2))

# Authentication functions
def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    return pwd_context.hash(password)

def get_user(username: str):
    users_db = get_users_db()
    if username in users_db:
        user_data = users_db[username].copy()
        # Remove hashed_password from user data before returning
        if "hashed_password" in user_data:
            user_data.pop("hashed_password")
        return User(**user_data)
    return None

def authenticate_user(username: str, password: str):
    users_db = get_users_db()
    if username not in users_db:
        return False
    user = get_user(username)
    if not user:
        return False
    if not verify_password(password, users_db[username]["hashed_password"]):
        return False
    return user

def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

async def get_current_user(token: str = Depends(oauth2_scheme)):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        tenant: str = payload.get("tenant")
        role: str = payload.get("role", "user")
        if username is None:
            raise credentials_exception
        token_data = TokenData(username=username, tenant=tenant, role=role)
    except JWTError:
        raise credentials_exception
    user = get_user(token_data.username)
    if user is None:
        raise credentials_exception
    return user

async def get_current_active_user(current_user: User = Depends(get_current_user)):
    users_db = get_users_db()
    if users_db.get(current_user.username, {}).get("disabled", False):
        raise HTTPException(status_code=400, detail="Inactive user")
    return current_user

async def get_admin_user(current_user: User = Depends(get_current_active_user)):
    if current_user.role != "admin":
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Not enough permissions"
        )
    return current_user

# User management functions
def create_user(user_data: UserCreate):
    users_db = get_users_db()
    if user_data.username in users_db:
        return False
    
    hashed_password = get_password_hash(user_data.password)
    users_db[user_data.username] = {
        "username": user_data.username,
        "tenant": user_data.tenant,
        "role": user_data.role,
        "disabled": user_data.disabled,
        "hashed_password": hashed_password
    }
    save_users_db(users_db)
    return True

def update_user(username: str, user_data: dict):
    users_db = get_users_db()
    if username not in users_db:
        return False
    
    for key, value in user_data.items():
        if key != "username" and key != "hashed_password":
            users_db[username][key] = value
    
    if "password" in user_data:
        users_db[username]["hashed_password"] = get_password_hash(user_data["password"])
    
    save_users_db(users_db)
    return True

def delete_user(username: str):
    users_db = get_users_db()
    if username not in users_db:
        return False
    
    del users_db[username]
    save_users_db(users_db)
    return True

# ───────────── Config & storage path helpers ───────────────
def cfg_path(t: str, a: str) -> Path: 
    return BASE_CONFIG_DIR / t / f"{a}.json"

def store_path(t: str, a: str) -> Path: 
    return BASE_STORE_DIR / t / a

def load_config(t: str, a: str) -> Dict[str, object]:
    p = cfg_path(t, a)
    if p.exists():
        return json.loads(p.read_text())
    p.parent.mkdir(parents=True, exist_ok=True)
    cfg = {
        "bot_name": f"{t}-{a}-Bot", 
        "system_prompt": "You are a helpful assistant.",
        "primary_color": "#1E88E5", 
        "secondary_color": "#FFFFFF", 
        "avatar_url": "",
        "mode": "inline", 
        "auto_open": False,
        "llm_provider": "openai",
        "llm_model": "gpt-4o-mini",
        "temperature": 0.3,
        "allowed_domains": ["*"]
    }
    p.write_text(json.dumps(cfg, indent=2))
    return cfg

def save_config(t: str, a: str, cfg: Dict[str, Any]) -> bool:
    p = cfg_path(t, a)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(cfg, indent=2))
    return True

# ─────────────────── Google Drive helpers ───────────────────
def _drive_service():
    creds_path = os.environ.get("GOOGLE_APPLICATION_CREDENTIALS")
    if not creds_path or not Path(creds_path).exists():
        raise RuntimeError("GOOGLE_APPLICATION_CREDENTIALS not set or file missing")
    creds = service_account.Credentials.from_service_account_file(
        creds_path,
        scopes=["https://www.googleapis.com/auth/drive.readonly"]
    )
    return build("drive", "v3", credentials=creds, cache_discovery=False)

def _drive_list(folder_id: str):
    q = f"'{folder_id}' in parents and trashed=false"
    return _drive_service().files().list(
        q=q, 
        fields="files(id,name)"
    ).execute().get("files", [])

@contextlib.contextmanager
def _drive_download(fid: str):
    request = _drive_service().files().get_media(fileId=fid)
    fd, path = tempfile.mkstemp()
    with os.fdopen(fd, "wb") as fh:
        dl = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            _, done = dl.next_chunk()
    try: 
        yield Path(path)
    finally: 
        os.remove(path)

# ───────────────────── Sitemap helpers ──────────────────────
def _parse_sitemap(url: str) -> List[str]:
    r = requests.get(url, timeout=15)
    r.raise_for_status()
    root = ET.fromstring(r.text)
    ns = {"sm": root.tag.split("}")[0].strip("{")}
    locs = root.findall("sm:url/sm:loc", ns) or root.findall("sm:sitemap/sm:loc", ns)
    urls = [e.text.strip() for e in locs]
    if root.find("sm:sitemap", ns) is not None:
        nested = []
        for sub in urls: 
            nested.extend(_parse_sitemap(sub))
        return nested
    return urls

def _download_page(u: str) -> str:
    r = requests.get(u, timeout=15, headers={"User-Agent": "CQBotCrawler"})
    r.raise_for_status()
    soup = BeautifulSoup(r.text, "html.parser")
    for t in soup(["script", "style", "noscript"]): 
        t.decompose()
    return soup.get_text("\n", strip=True)

# ───────────────────────── Ingest ───────────────────────────
def ingest(t: str, a: str, sitemap: Optional[str] = None, drive: Optional[str] = None, 
          files: Optional[List[Path]] = None, console: Optional[Console] = None):
    """
    Ingest content from various sources into the vector store.
    
    Args:
        t: Tenant identifier
        a: Agent identifier
        sitemap: Optional URL to sitemap
        drive: Optional Google Drive folder ID
        files: Optional list of local files to process
        console: Optional Rich console for progress display
    """
    txts, metas = [], []
    emb = OpenAIEmbeddings()
    
    if console:
        console.print(Panel.fit(f"Starting ingestion for tenant [bold]{t}[/bold], agent [bold]{a}[/bold]"))
    
    # Process Google Drive files
    if drive:
        files_list = _drive_list(drive)
        if console:
            with Progress() as progress:
                task = progress.add_task(f"Processing [bold]{len(files_list)}[/bold] Drive files...", total=len(files_list))
                for f in files_list:
                    progress.update(task, advance=1, description=f"Processing {f['name']}")
                    with _drive_download(f["id"]) as p:
                        _process_file(p, f["name"], txts, metas)
        else:
            for f in files_list:
                with _drive_download(f["id"]) as p:
                    _process_file(p, f["name"], txts, metas)
    
    # Process local files
    if files:
        if console:
            with Progress() as progress:
                task = progress.add_task(f"Processing [bold]{len(files)}[/bold] local files...", total=len(files))
                for file in files:
                    progress.update(task, advance=1, description=f"Processing {file.name}")
                    _process_file(file, file.name, txts, metas)
        else:
            for file in files:
                _process_file(file, file.name, txts, metas)
    
    # Process sitemap URLs
    if sitemap:
        try:
            urls = _parse_sitemap(sitemap)
            if console:
                console.print(f"Found [bold]{len(urls)}[/bold] URLs in sitemap")
                with Progress() as progress:
                    task = progress.add_task(f"Processing URLs...", total=len(urls))
                    for u in urls:
                        try:
                            progress.update(task, advance=1, description=f"Processing {u}")
                            pg = _download_page(u)
                            for c in TEXT_SPLITTER.split_text(pg):
                                txts.append(c)
                                metas.append({"source": u})
                        except Exception as e:
                            progress.console.print(f"Error processing {u}: {str(e)}", style="red")
            else:
                for u in urls:
                    try:
                        pg = _download_page(u)
                        for c in TEXT_SPLITTER.split_text(pg):
                            txts.append(c)
                            metas.append({"source": u})
                    except Exception:
                        pass
        except Exception as e:
            if console:
                console.print(f"[red]Error processing sitemap: {str(e)}[/red]")
            else:
                print(f"Error processing sitemap: {str(e)}")

    if not txts:
        msg = "Nothing to ingest"
        if console:
            console.print(f"[yellow]{msg}[/yellow]")
        else:
            print(msg)
        return

    # Save to vector store
    path = store_path(t, a)
    path.mkdir(parents=True, exist_ok=True)
    
    if console:
        console.print(f"Processing [bold]{len(txts)}[/bold] text chunks into vector store...")
        
    FAISS.from_texts(txts, emb, metadatas=metas).save_local(str(path))
    
    msg = f"Vector store saved to {path}"
    if console:
        console.print(f"[green]{msg}[/green]")
    else:
        print(msg)

def _process_file(file_path: Path, filename: str, txts: List[str], metas: List[Dict]):
    """Process a file and extract text chunks for vectorization"""
    ext = file_path.suffix.lower()
    raw = ""
    
    try:
        if ext == ".pdf":
            import pypdf
            raw = "\n".join(pg.extract_text() or "" for pg in pypdf.PdfReader(str(file_path)).pages)
        elif ext in {".txt", ".md"}:
            raw = file_path.read_text(errors='ignore')
        elif ext == ".docx":
            import docx
            doc = docx.Document(str(file_path))
            raw = "\n".join(para.text for para in doc.paragraphs)
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            raw = "\n".join(shape.text for slide in prs.slides 
                          for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path))
            raw = "\n".join(
                "\n".join(str(cell.value) for cell in row if cell.value)
                for sheet in wb.worksheets for row in sheet.rows
            )
        elif ext in {".csv"}:
            import csv
            with open(str(file_path), 'r', errors='ignore') as csvfile:
                reader = csv.reader(csvfile)
                raw = "\n".join(",".join(str(cell) for cell in row) for row in reader)
        elif ext in {".html", ".htm"}:
            with open(str(file_path), 'r', errors='ignore') as htmlfile:
                soup = BeautifulSoup(htmlfile.read(), "html.parser")
                raw = soup.get_text("\n", strip=True)
        else:
            return
            
        for c in TEXT_SPLITTER.split_text(raw):
            txts.append(c)
            metas.append({"source": filename})
    except Exception as e:
        print(f"Error processing file {filename}: {str(e)}")

# ────────────────── LLM Provider Selection ─────────────────
def get_llm_response(messages: List[Dict], provider="openai", model=None, temperature=0.3):
    """Get response from selected LLM provider"""
    start_time = time.time()
    tokens_in = _estimate_tokens(messages)
    
    try:
        if provider == "openai":
            import openai
            openai.api_key = os.getenv("OPENAI_API_KEY", "")
            model = model or "gpt-4o-mini"
            rsp = openai.ChatCompletion.create(
                model=model, 
                temperature=temperature, 
                messages=messages
            )
            content = rsp.choices[0].message["content"]
            tokens_out = rsp.usage.completion_tokens
            
        elif provider == "anthropic":
            import anthropic
            client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY", ""))
            model = model or "claude-3-opus-20240229"
            rsp = client.messages.create(
                model=model,
                max_tokens=1000,
                temperature=temperature,
                messages=[
                    {"role": m["role"], "content": m["content"]} for m in messages
                ]
            )
            content = rsp.content[0].text
            tokens_out = _estimate_tokens([{"role": "assistant", "content": content}])
            
        elif provider == "vertexai":
            from vertexai.generative_models import GenerativeModel
            model_name = model or "gemini-1.5-pro"
            model = GenerativeModel(model_name)
            response = model.generate_content([
                {"role": m["role"], "parts": [{"text": m["content"]}]} 
                for m in messages
            ])
            content = response.text
            tokens_out = _estimate_tokens([{"role": "assistant", "content": content}])
            
        else:
            raise ValueError(f"Unknown LLM provider: {provider}")
            
    except Exception as e:
        content = f"Error generating response: {str(e)}"
        tokens_out = len(content) // 4  # rough estimate
    
    latency = time.time() - start_time
    
    return {
        "content": content,
        "latency": latency,
        "tokens_in": tokens_in,
        "tokens_out": tokens_out
    }

def _estimate_tokens(messages):
    """Roughly estimate token count based on character count"""
    total_chars = sum(len(m.get("content", "")) for m in messages)
    return total_chars // 4  # rough estimate: 4 chars per token

# ───────────────────────── API ──────────────────────────────
app = FastAPI(title="Multi‑Tenant RAG", version="7.0")
app.add_middleware(
    CORSMiddleware, 
    allow_origins=["*"], 
    allow_methods=["*"], 
    allow_headers=["*"]
)

# API models
class ChatRequest(BaseModel):
    messages: List[dict]
    
class ChatResponse(BaseModel):
    reply: str
    sources: List[dict]
    
class ConfigUpdateRequest(BaseModel):
    bot_name: Optional[str] = None
    system_prompt: Optional[str] = None
    primary_color: Optional[str] = None
    secondary_color: Optional[str] = None
    avatar_url: Optional[str] = None
    mode: Optional[str] = None
    auto_open: Optional[bool] = None
    llm_provider: Optional[str] = None
    llm_model: Optional[str] = None
    temperature: Optional[float] = None
    allowed_domains: Optional[List[str]] = None
    
# Vector store cache
_vec: Dict[str, FAISS] = {}

def _db(t: str, a: str):
    k = f"{t}/{a}"
    p = store_path(t, a)
    if k in _vec: 
        return _vec[k]
    if not p.exists(): 
        raise HTTPException(404, "Vector store missing; run ingest")
    _vec[k] = FAISS.load_local(str(p), OpenAIEmbeddings())
    return _vec[k]

# ─────────────────── Authentication Endpoints ───────────────
@app.post("/token", response_model=Token)
async def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends()):
    user = authenticate_user(form_data.username, form_data.password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    access_token = create_access_token(
        data={"sub": user.username, "tenant": user.tenant, "role": user.role}
    )
    return {"access_token": access_token, "token_type": "bearer"}

@app.get("/users/me", response_model=User)
async def read_users_me(current_user: User = Depends(get_current_active_user)):
    return current_user

@app.post("/users", status_code=status.HTTP_201_CREATED)
async def create_new_user(user: UserCreate, admin: User = Depends(get_admin_user)):
    if create_user(user):
        return {"message": "User created successfully"}
    raise HTTPException(
        status_code=status.HTTP_400_BAD_REQUEST,
        detail="Username already exists"
    )

@app.put("/users/{username}")
async def update_existing_user(username: str, user_data: dict, admin: User = Depends(get_admin_user)):
    if update_user(username, user_data):
        return {"message": "User updated successfully"}
    raise HTTPException(
        status_code=status.HTTP_404_NOT_FOUND,
        detail="User not found"
    )

@app.delete("/users/{username}")
async def delete_existing_user(username: str, admin: User = Depends(get_admin_user)):
    if delete_user(username):
        return {"message": "User deleted successfully"}
    raise HTTPException(
        status_code=status.HTTP_404_NOT_FOUND,
        detail="User not found"
    )

@app.get("/users")
async def list_users(admin: User = Depends(get_admin_user)):
    users_db = get_users_db()
    # Remove sensitive information
    users = []
    for username, user_data in users_db.items():
        user_info = user_data.copy()
        if "hashed_password" in user_info:
            user_info.pop("hashed_password")
        users.append(user_info)
    return users

# ─────────────────── Chat Endpoints ───────────────────────
@app.post("/chat", response_model=ChatResponse)
async def chat(
    req: ChatRequest, 
    request: Request, 
    tenant: str = Query(DEFAULT_TENANT), 
    agent: str = Query(DEFAULT_AGENT),
    current_user: User = Depends(get_current_active_user)
):
    # Check if user has access to this tenant
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant"
        )
    
    # Get vector database and configuration
    db = _db(tenant, agent)
    cfg = load_config(tenant, agent)
    
    # Get the latest user question
    q = next((m["content"] for m in reversed(req.messages) if m["role"] == "user"), "")
    
    # Retrieve relevant documents
    docs = db.similarity_search_with_score(q, k=4)
    ctx = "\n".join(d.page_content for d, _ in docs)
    
    # Create system message with context
    system_msg = {
        "role": "system", 
        "content": cfg["system_prompt"] + "\nContext:\n" + ctx
    }
    
    # Get response from LLM
    llm_result = get_llm_response(
        messages=[system_msg, *req.messages],
        provider=cfg.get("llm_provider", "openai"),
        model=cfg.get("llm_model", "gpt-4o-mini"),
        temperature=cfg.get("temperature", 0.3)
    )
    
    # Extract sources
    srcs, seen = [], set()
    for d, _ in docs:
        s = d.metadata.get("source", "")
        if s and s not in seen:
            srcs.append({"source": s})
            seen.add(s)
    
    # Log the interaction
    with sqlite3.connect(DB_PATH) as con:
        con.execute(
            """INSERT INTO chat_logs
               (ts, tenant, agent, session_id, question, answer, sources, 
                latency, tokens_in, tokens_out, user_ip)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (
                datetime.now(timezone.utc).isoformat(),
                tenant,
                agent,
                request.headers.get("X-Session-Id", "anon"),
                q,
                llm_result["content"],
                json.dumps(srcs),
                llm_result["latency"],
                llm_result["tokens_in"],
                llm_result["tokens_out"],
                request.client.host
            )
        )
    
    return {
        "reply": llm_result["content"],
        "sources": srcs
    }

@app.post("/feedback/{chat_id}")
async def submit_feedback(
    chat_id: int,
    feedback: int,
    current_user: User = Depends(get_current_active_user)
):
    # Validate feedback score (1-5)
    if feedback < 1 or feedback > 5:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Feedback must be between 1 and 5"
        )
    
    # Update the feedback in the database
    with sqlite3.connect(DB_PATH) as con:
        result = con.execute(
            "UPDATE chat_logs SET user_feedback = ? WHERE id = ?",
            (feedback, chat_id)
        )
        if result.rowcount == 0:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Chat log not found"
            )
    
    return {"message": "Feedback submitted successfully"}

# ────────────────── Config Endpoints ──────────────────────
@app.get("/config")
async def get_config(
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT),
    request: Request = None
):
    cfg = load_config(tenant, agent)
    
    # Check if domain is allowed (for widget embedding)
    if request and "Origin" in request.headers:
        origin = request.headers["Origin"]
        allowed_domains = cfg.get("allowed_domains", ["*"])
        
        if "*" not in allowed_domains:
            from urllib.parse import urlparse
            domain = urlparse(origin).netloc
            if domain not in allowed_domains:
                raise HTTPException(
                    status_code=status.HTTP_403_FORBIDDEN,
                    detail="Domain not allowed to access this configuration"
                )
    
    # Remove sensitive or internal configuration
    public_cfg = {k: v for k, v in cfg.items() if k not in [
        "llm_provider", "llm_model", "temperature", "allowed_domains"
    ]}
    
    return public_cfg

@app.put("/config")
async def update_config(
    config: ConfigUpdateRequest,
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT),
    current_user: User = Depends(get_admin_user)
):
    cfg = load_config(tenant, agent)
    
    # Update only provided fields
    for field, value in config.dict(exclude_unset=True).items():
        if value is not None:
            cfg[field] = value
    
    save_config(tenant, agent, cfg)
    return {"message": "Configuration updated successfully"}

@app.get("/tenants")
async def list_tenants(current_user: User = Depends(get_admin_user)):
    tenants = []
    for tenant_dir in BASE_CONFIG_DIR.iterdir():
        if tenant_dir.is_dir():
            agents = []
            for config_file in tenant_dir.iterdir():
                if config_file.is_file() and config_file.suffix == ".json":
                    agent_name = config_file.stem
                    agents.append(agent_name)
            tenants.append({"tenant": tenant_dir.name, "agents": agents})
    return tenants

# ─────────────────── Analytics Endpoints ───────────────────
@app.get("/analytics")
async def get_analytics(
    tenant: str = Query(DEFAULT_TENANT),
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    current_user: User = Depends(get_admin_user)
):
    # Check if user has access to this tenant
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant's analytics"
        )
    
    # Set default date range if not provided (last 7 days)
    if not end_date:
        end_date = datetime.now(timezone.utc).isoformat()
    if not start_date:
        start_date = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
    
    # Query the database
    with sqlite3.connect(DB_PATH) as con:
        df = pd.read_sql_query(
            "SELECT * FROM chat_logs WHERE tenant = ? AND ts BETWEEN ? AND ?",
            con, params=(tenant, start_date, end_date)
        )
    
    # Generate analytics
    if df.empty:
        return {"message": "No data available for selected period"}
    
    # Daily usage counts
    df['date'] = pd.to_datetime(df['ts']).dt.date
    daily_counts = df.groupby(['date']).size().reset_index(name='count')
    
    # Top questions
    question_counts = df['question'].value_counts().head(10).to_dict()
    
    # Source utilization
    sources = []
    for src_json in df['sources'].dropna():
        try:
            src_list = json.loads(src_json)
            for src in src_list:
                if 'source' in src:
                    sources.append(src['source'])
        except:
            pass
    source_counts = pd.Series(sources).value_counts().head(10).to_dict()
    
    # Session statistics
    session_counts = df.groupby('session_id').size().describe().to_dict()
    
    # Performance metrics
    latency_stats = df['latency'].describe().to_dict()
    token_stats = {
        'input': df['tokens_in'].describe().to_dict(),
        'output': df['tokens_out'].describe().to_dict()
    }
    
    # Feedback stats
    feedback_counts = df['user_feedback'].value_counts().to_dict()
    avg_feedback = df['user_feedback'].mean() if 'user_feedback' in df else None
    
    return {
        "total_queries": len(df),
        "unique_sessions": df['session_id'].nunique(),
        "daily_counts": daily_counts.to_dict('records'),
        "top_questions": question_counts,
        "top_sources": source_counts,
        "session_stats": session_counts,
        "performance": latency_stats,
        "token_stats": token_stats,
        "feedback": {
            "counts": feedback_counts,
            "average": avg_feedback
        }
    }

# ─────────────────── Widget Endpoint ───────────────────────
@app.get("/widget.js", response_class=PlainTextResponse)
async def get_widget():
    return WIDGET_JS

# Complete JavaScript widget with advanced features
WIDGET_JS = r"""(function(){
const p=new URLSearchParams(location.search);
const tenant=p.get('tenant')||'public';
const agent=p.get('agent')||'default';
const sid=sessionStorage.getItem('cq_sid')||(()=>{const r=Math.random().toString(36).slice(2);sessionStorage.setItem('cq_sid',r);return r})();
const msgs=[];
function $(id){return document.getElementById(id);}

// Create widget container
const container = document.createElement('div');
container.id = 'cq-widget-container';
document.body.appendChild(container);

// Advanced features
let config = {};
const features = {
  typing: true,        // Show typing indicator
  fileAttachments: true, // Allow file attachments
  voiceInput: true,    // Voice input capability
  darkMode: false,     // Dark mode toggle
  responseSpeech: false // Text-to-speech for responses
};

// Authentication mechanism
let authToken = localStorage.getItem('cq_auth_token');
const headers = {
  'Content-Type': 'application/json',
  'X-Session-Id': sid
};
if (authToken) headers['Authorization'] = `Bearer ${authToken}`;

// Fetch configuration
fetch(`/config?tenant=${tenant}&agent=${agent}`)
  .then(r => r.json())
  .then(cfg => {
    config = cfg;
    initWidget(cfg);
  });

function initWidget(cfg) {
  // Create chat interface
  container.innerHTML = `
    <div id="cq-chat-widget" style="--primary:${cfg.primary_color};--secondary:${cfg.secondary_color}">
      <div id="cq-header">
        <img src="${cfg.avatar_url || 'default-avatar.png'}" alt="${cfg.bot_name}">
        <h3>${cfg.bot_name}</h3>
        <div class="cq-controls">
          <button id="cq-settings">⚙️</button>
          <button id="cq-close">×</button>
        </div>
      </div>
      <div id="cq-messages"></div>
      <div id="cq-typing" style="display:none">
        <div class="cq-dot"></div>
        <div class="cq-dot"></div>
        <div class="cq-dot"></div>
      </div>
      <div id="cq-footer">
        <div id="cq-attachments"></div>
        <textarea id="cq-input" placeholder="Type your message..."></textarea>
        <div class="cq-actions">
          <button id="cq-mic" title="Voice input">🎤</button>
          <button id="cq-attach" title="Attach file">📎</button>
          <button id="cq-send">➤</button>
        </div>
      </div>
      <div id="cq-settings-panel" style="display:none">
        <h4>Settings</h4>
        <label><input type="checkbox" id="cq-dark-mode"> Dark mode</label>
        <label><input type="checkbox" id="cq-speech"> Read responses aloud</label>
        <h4>Sources</h4>
        <div id="cq-sources"></div>
      </div>
    </div>
  `;
  
  // Set up event handlers
  $('cq-send').addEventListener('click', sendMessage);
  $('cq-input').addEventListener('keypress', e => e.key === 'Enter' && !e.shiftKey && sendMessage());
  $('cq-settings').addEventListener('click', toggleSettings);
  $('cq-close').addEventListener('click', minimizeWidget);
  $('cq-mic').addEventListener('click', startVoiceInput);
  $('cq-attach').addEventListener('click', () => $('cq-file-input').click());
  $('cq-dark-mode').addEventListener('change', e => setDarkMode(e.target.checked));
  
  // If auto-open is configured, show widget
  if (cfg.auto_open) showWidget();
}

function sendMessage() {
  const input = $('cq-input');
  const text = input.value.trim();
  if (!text) return;
  
  addMessage('user', text);
  input.value = '';
  
  // Show typing indicator
  if (features.typing) $('cq-typing').style.display = 'flex';
  
  // Send to backend with authentication
  fetch(`/chat?tenant=${tenant}&agent=${agent}`, {
    method: 'POST',
    headers: headers,
    body: JSON.stringify({messages: msgs})
  })
  .then(r => r.json())
  .then(data => {
    if (features.typing) $('cq-typing').style.display = 'none';
    addMessage('assistant', data.reply);
    updateSources(data.sources);
    
    // Text-to-speech if enabled
    if (features.responseSpeech) {
      const speech = new SpeechSynthesisUtterance(data.reply);
      window.speechSynthesis.speak(speech);
    }
  });
}

// Additional methods for handling files, voice input, etc.
function startVoiceInput() {
  if (!('webkitSpeechRecognition' in window)) {
    alert('Voice input not supported in your browser');
    return;
  }
  
  const recognition = new webkitSpeechRecognition();
  recognition.continuous = false;
  recognition.interimResults = false;
  
  recognition.onresult = function(event) {
    const transcript = event.results[0][0].transcript;
    $('cq-input').value = transcript;
  };
  
  recognition.start();
}

function addMessage(role, content) {
  msgs.push({role, content});
  
  const msgDiv = document.createElement('div');
  msgDiv.className = `cq-message cq-${role}`;
  msgDiv.innerHTML = `
    <div class="cq-avatar">${role === 'user' ? '👤' : '🤖'}</div>
    <div class="cq-bubble">${formatMessage(content)}</div>
  `;
  
  $('cq-messages').appendChild(msgDiv);
  $('cq-messages').scrollTop = $('cq-messages').scrollHeight;
}

function formatMessage(text) {
  // Simple markdown-like formatting
  return text
    .replace(/\*\*(.*?)\*\*/g, '<strong>$1</strong>')
    .replace(/\*(.*?)\*/g, '<em>$1</em>')
    .replace(/```(.*?)```/gs, '<pre><code>$1</code></pre>')
    .replace(/`(.*?)`/g, '<code>$1</code>')
    .replace(/\n/g, '<br>');
}

function updateSources(sources) {
  const sourcesDiv = $('cq-sources');
  sourcesDiv.innerHTML = '';
  
  if (sources && sources.length > 0) {
    const sourceList = document.createElement('ul');
    sources.forEach(src => {
      const li = document.createElement('li');
      li.innerHTML = `<a href="${src.source}" target="_blank">${src.source}</a>`;
      sourceList.appendChild(li);
    });
    sourcesDiv.appendChild(sourceList);
  }
}

function toggleSettings() {
  const panel = $('cq-settings-panel');
  panel.style.display = panel.style.display === 'none' ? 'block' : 'none';
}

function minimizeWidget() {
  const widget = $('cq-chat-widget');
  widget.style.display = 'none';
  $('cq-launcher').style.display = 'block';
}

function showWidget() {
  const widget = $('cq-chat-widget');
  widget.style.display = 'flex';
  $('cq-launcher').style.display = 'none';
}

function toggleWidget() {
  const widget = $('cq-chat-widget');
  if (widget.style.display === 'none' || !widget.style.display) {
    showWidget();
  } else {
    minimizeWidget();
  }
}

function setDarkMode(enabled) {
  features.darkMode = enabled;
  const widget = $('cq-chat-widget');
  if (enabled) {
    widget.classList.add('cq-dark');
  } else {
    widget.classList.remove('cq-dark');
  }
  localStorage.setItem('cq_dark_mode', enabled);
}

// Initialize widget style
const style = document.createElement('style');
style.textContent = `
  #cq-widget-container {
    position: fixed;
    bottom: 20px;
    right: 20px;
    z-index: 10000;
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
  }
  
  #cq-launcher {
    width: 60px;
    height: 60px;
    border-radius: 50%;
    background: var(--primary, #1E88E5);
    color: white;
    border: none;
    font-size: 24px;
    cursor: pointer;
    box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    transition: all 0.3s ease;
  }
  
  #cq-launcher:hover {
    transform: scale(1.1);
    box-shadow: 0 6px 16px rgba(0,0,0,0.2);
  }
  
  #cq-chat-widget {
    width: 350px;
    height: 500px;
    border-radius: 12px;
    overflow: hidden;
    display: flex;
    flex-direction: column;
    box-shadow: 0 8px 32px rgba(0,0,0,0.12);
    background: white;
    transition: all 0.3s ease;
    border: 1px solid #e0e0e0;
  }
  
  #cq-chat-widget.cq-dark {
    background: #2d2d2d;
    border-color: #404040;
    color: white;
  }
  
  #cq-header {
    background: var(--primary, #1E88E5);
    color: white;
    padding: 16px;
    display: flex;
    align-items: center;
    gap: 12px;
  }
  
  #cq-header img {
    width: 32px;
    height: 32px;
    border-radius: 50%;
    object-fit: cover;
  }
  
  #cq-header h3 {
    margin: 0;
    flex: 1;
    font-size: 16px;
    font-weight: 600;
  }
  
  .cq-controls {
    display: flex;
    gap: 8px;
  }
  
  .cq-controls button {
    background: rgba(255,255,255,0.2);
    border: none;
    color: white;
    width: 32px;
    height: 32px;
    border-radius: 6px;
    cursor: pointer;
    transition: background 0.2s;
  }
  
  .cq-controls button:hover {
    background: rgba(255,255,255,0.3);
  }
  
  #cq-messages {
    flex: 1;
    padding: 16px;
    overflow-y: auto;
    display: flex;
    flex-direction: column;
    gap: 12px;
  }
  
  .cq-message {
    display: flex;
    gap: 8px;
    align-items: flex-start;
  }
  
  .cq-message.cq-user {
    flex-direction: row-reverse;
  }
  
  .cq-avatar {
    width: 32px;
    height: 32px;
    border-radius: 50%;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 14px;
    flex-shrink: 0;
  }
  
  .cq-bubble {
    max-width: 80%;
    padding: 12px 16px;
    border-radius: 18px;
    line-height: 1.4;
    word-wrap: break-word;
  }
  
  .cq-user .cq-bubble {
    background: var(--primary, #1E88E5);
    color: white;
    margin-left: auto;
  }
  
  .cq-assistant .cq-bubble {
    background: #f5f5f5;
    color: #333;
  }
  
  .cq-dark .cq-assistant .cq-bubble {
    background: #404040;
    color: white;
  }
  
  #cq-typing {
    padding: 16px;
    display: flex;
    justify-content: center;
    gap: 4px;
  }
  
  .cq-dot {
    width: 8px;
    height: 8px;
    border-radius: 50%;
    background: #ccc;
    animation: cq-bounce 1.4s infinite ease-in-out both;
  }
  
  .cq-dot:nth-child(1) { animation-delay: -0.32s; }
  .cq-dot:nth-child(2) { animation-delay: -0.16s; }
  
  @keyframes cq-bounce {
    0%, 80%, 100% { transform: scale(0); }
    40% { transform: scale(1); }
  }
  
  #cq-footer {
    padding: 16px;
    border-top: 1px solid #e0e0e0;
  }
  
  .cq-dark #cq-footer {
    border-top-color: #404040;
  }
  
  #cq-input {
    width: 100%;
    min-height: 20px;
    max-height: 100px;
    padding: 12px;
    border: 1px solid #e0e0e0;
    border-radius: 24px;
    resize: none;
    outline: none;
    font-family: inherit;
    font-size: 14px;
    line-height: 1.4;
  }
  
  .cq-dark #cq-input {
    background: #404040;
    border-color: #555;
    color: white;
  }
  
  .cq-actions {
    display: flex;
    justify-content: space-between;
    align-items: center;
    margin-top: 8px;
  }
  
  .cq-actions button {
    background: none;
    border: none;
    cursor: pointer;
    padding: 8px;
    border-radius: 50%;
    transition: background 0.2s;
  }
  
  .cq-actions button:hover {
    background: rgba(0,0,0,0.1);
  }
  
  .cq-dark .cq-actions button:hover {
    background: rgba(255,255,255,0.1);
  }
  
  #cq-send {
    background: var(--primary, #1E88E5) !important;
    color: white !important;
  }
  
  #cq-settings-panel {
    position: absolute;
    top: 60px;
    right: 0;
    background: white;
    border: 1px solid #e0e0e0;
    border-radius: 8px;
    padding: 16px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    min-width: 200px;
  }
  
  .cq-dark #cq-settings-panel {
    background: #2d2d2d;
    border-color: #404040;
    color: white;
  }
  
  #cq-settings-panel h4 {
    margin: 0 0 12px 0;
    font-size: 14px;
    font-weight: 600;
  }
  
  #cq-settings-panel label {
    display: flex;
    align-items: center;
    gap: 8px;
    margin-bottom: 8px;
    cursor: pointer;
  }
  
  #cq-sources ul {
    list-style: none;
    padding: 0;
    margin: 8px 0 0 0;
  }
  
  #cq-sources li {
    margin-bottom: 4px;
  }
  
  #cq-sources a {
    color: var(--primary, #1E88E5);
    text-decoration: none;
    font-size: 12px;
  }
  
  #cq-sources a:hover {
    text-decoration: underline;
  }
  
  .cq-dark #cq-sources a {
    color: #4fc3f7;
  }
  
  /* Mobile responsive */
  @media (max-width: 480px) {
    #cq-chat-widget {
      width: calc(100vw - 40px);
      height: calc(100vh - 40px);
      max-height: 600px;
    }
  }
`;
document.head.appendChild(style);

// Initialize if DOM already loaded, otherwise wait
if (document.readyState === 'loading') {
  document.addEventListener('DOMContentLoaded', init);
} else {
  init();
}

function init() {
  // Create launcher button
  const launcher = document.createElement('button');
  launcher.id = 'cq-launcher';
  launcher.innerHTML = '💬';
  launcher.onclick = toggleWidget;
  launcher.style.display = 'block';
  container.appendChild(launcher);
  
  // Load saved preferences
  const savedDarkMode = localStorage.getItem('cq_dark_mode') === 'true';
  if (savedDarkMode) {
    features.darkMode = true;
  }
}

})();"""

# ─────────────────── File Upload Endpoint ───────────────────
@app.post("/upload")
async def upload_files(
    files: List[UploadFile] = File(...),
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT),
    current_user: User = Depends(get_admin_user)
):
    """Upload and ingest files into the vector store"""
    
    # Check tenant access
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant"
        )
    
    temp_files = []
    try:
        # Save uploaded files temporarily
        for file in files:
            temp_path = Path(tempfile.mkdtemp()) / file.filename
            temp_path.parent.mkdir(exist_ok=True)
            
            with temp_path.open("wb") as buffer:
                content = await file.read()
                buffer.write(content)
            
            temp_files.append(temp_path)
        
        # Ingest the files
        ingest(tenant, agent, files=temp_files)
        
        # Clear vector store cache to force reload
        cache_key = f"{tenant}/{agent}"
        if cache_key in _vec:
            del _vec[cache_key]
        
        return {
            "message": f"Successfully uploaded and processed {len(files)} files",
            "files": [f.filename for f in files]
        }
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error processing files: {str(e)}"
        )
    finally:
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                temp_file.unlink()
                temp_file.parent.rmdir()
            except:
                pass

# ─────────────────── Ingest Endpoints ───────────────────────
@app.post("/ingest/sitemap")
async def ingest_sitemap(
    sitemap_url: str,
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT),
    current_user: User = Depends(get_admin_user)
):
    """Ingest content from a sitemap URL"""
    
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant"
        )
    
    try:
        ingest(tenant, agent, sitemap=sitemap_url)
        
        # Clear cache
        cache_key = f"{tenant}/{agent}"
        if cache_key in _vec:
            del _vec[cache_key]
        
        return {"message": f"Successfully ingested content from sitemap: {sitemap_url}"}
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error ingesting sitemap: {str(e)}"
        )

@app.post("/ingest/drive")
async def ingest_drive(
    folder_id: str,
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT),
    current_user: User = Depends(get_admin_user)
):
    """Ingest content from a Google Drive folder"""
    
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant"
        )
    
    try:
        ingest(tenant, agent, drive=folder_id)
        
        # Clear cache
        cache_key = f"{tenant}/{agent}"
        if cache_key in _vec:
            del _vec[cache_key]
        
        return {"message": f"Successfully ingested content from Google Drive folder: {folder_id}"}
        
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error ingesting Google Drive content: {str(e)}"
        )

# ─────────────────── CLI with GUI ───────────────────────────
cli_app = typer.Typer(name="rag-chatbot", help="Multi-tenant RAG Chatbot CLI")
console = Console()

@cli_app.command("dashboard")
def dashboard():
    """Launch the interactive dashboard"""
    
    def show_main_menu():
        console.clear()
        console.print(Panel.fit(
            "[bold blue]RAG Chatbot Administration Dashboard[/bold blue]",
            title="Welcome",
            border_style="blue"
        ))
        
        console.print("\n[bold]Available Commands:[/bold]")
        console.print("1. 📊 View Analytics")
        console.print("2. 🏢 Manage Tenants & Agents")
        console.print("3. 👥 Manage Users")
        console.print("4. 📁 Ingest Content")
        console.print("5. ⚙️  System Status")
        console.print("6. 🚪 Exit")
        
        choice = typer.prompt("\nSelect an option")
        return choice
    
    def show_analytics():
        console.clear()
        console.print(Panel.fit("Analytics Dashboard", border_style="green"))
        
        # Get available tenants
        tenants = []
        for tenant_dir in BASE_CONFIG_DIR.iterdir():
            if tenant_dir.is_dir():
                tenants.append(tenant_dir.name)
        
        if not tenants:
            console.print("[yellow]No tenants found[/yellow]")
            typer.prompt("Press Enter to continue")
            return
        
        # Show tenant analytics
        for tenant in tenants:
            with sqlite3.connect(DB_PATH) as con:
                cursor = con.execute(
                    "SELECT COUNT(*) as total, COUNT(DISTINCT session_id) as sessions FROM chat_logs WHERE tenant = ?",
                    (tenant,)
                )
                total, sessions = cursor.fetchone()
            
            table = Table(title=f"Tenant: {tenant}")
            table.add_column("Metric")
            table.add_column("Value")
            table.add_row("Total Queries", str(total))
            table.add_row("Unique Sessions", str(sessions))
            console.print(table)
        
        typer.prompt("\nPress Enter to continue")
    
    def show_tenants():
        console.clear()
        console.print(Panel.fit("Tenant & Agent Management", border_style="yellow"))
        
        table = Table()
        table.add_column("Tenant")
        table.add_column("Agents")
        table.add_column("Vector Store")
        
        for tenant_dir in BASE_CONFIG_DIR.iterdir():
            if tenant_dir.is_dir():
                agents = []
                for config_file in tenant_dir.iterdir():
                    if config_file.is_file() and config_file.suffix == ".json":
                        agents.append(config_file.stem)
                
                store_exists = "✅" if store_path(tenant_dir.name, agents[0] if agents else "default").exists() else "❌"
                table.add_row(tenant_dir.name, ", ".join(agents), store_exists)
        
        console.print(table)
        typer.prompt("\nPress Enter to continue")
    
   
def show_users():
    while True:
        console.clear()
        console.print(Panel.fit("User Management", border_style="red"))
        
        users_db = get_users_db()
        
        table = Table()
        table.add_column("Username")
        table.add_column("Tenant")
        table.add_column("Role")
        table.add_column("Status")
        
        for username, user_data in users_db.items():
            status = "❌ Disabled" if user_data.get("disabled", False) else "✅ Active"
            table.add_row(
                username,
                user_data.get("tenant", "N/A"),
                user_data.get("role", "user"),
                status
            )
        
        console.print(table)
        
        console.print("\n[bold]User Management Options:[/bold]")
        console.print("1. ➕ Create New User")
        console.print("2. ✏️  Edit User")
        console.print("3. ❌ Delete User")
        console.print("4. 🔄 Enable/Disable User")
        console.print("5. 🔙 Back to Main Menu")
        
        choice = typer.prompt("\nSelect an option")
        
        if choice == "1":
            # Create new user
            console.print("\n[bold green]Create New User[/bold green]")
            username = typer.prompt("Enter username")
            
            # Check if username already exists
            if username in users_db:
                console.print(f"[red]❌ User '{username}' already exists[/red]")
                typer.prompt("Press Enter to continue")
                continue
            
            password = typer.prompt("Enter password", hide_input=True)
            confirm_password = typer.prompt("Confirm password", hide_input=True)
            
            if password != confirm_password:
                console.print("[red]❌ Passwords don't match[/red]")
                typer.prompt("Press Enter to continue")
                continue
            
            # Show available tenants
            available_tenants = ["*"]  # Admin access to all
            for tenant_dir in BASE_CONFIG_DIR.iterdir():
                if tenant_dir.is_dir():
                    available_tenants.append(tenant_dir.name)
            
            console.print(f"\nAvailable tenants: {', '.join(available_tenants)}")
            tenant = typer.prompt("Enter tenant (* for all tenants)")
            
            console.print("\nAvailable roles: admin, user")
            role = typer.prompt("Enter role", default="user")
            
            if role not in ["admin", "user"]:
                console.print("[red]❌ Invalid role. Must be 'admin' or 'user'[/red]")
                typer.prompt("Press Enter to continue")
                continue
            
            # Create user
            user_data = UserCreate(
                username=username,
                password=password,
                tenant=tenant,
                role=role,
                disabled=False
            )
            
            if create_user(user_data):
                console.print(f"[green]✅ User '{username}' created successfully[/green]")
            else:
                console.print(f"[red]❌ Failed to create user '{username}'[/red]")
            
            typer.prompt("Press Enter to continue")
        
        elif choice == "2":
            # Edit user
            console.print("\n[bold yellow]Edit User[/bold yellow]")
            if not users_db:
                console.print("[yellow]No users found[/yellow]")
                typer.prompt("Press Enter to continue")
                continue
            
            username = typer.prompt(f"Enter username to edit ({', '.join(users_db.keys())})")
            
            if username not in users_db:
                console.print(f"[red]❌ User '{username}' not found[/red]")
                typer.prompt("Press Enter to continue")
                continue
            
            user_info = users_db[username]
            console.print(f"\nCurrent user info for '{username}':")
            console.print(f"Tenant: {user_info.get('tenant', 'N/A')}")
            console.print(f"Role: {user_info.get('role', 'user')}")
            console.print(f"Status: {'Disabled' if user_info.get('disabled', False) else 'Active'}")
            
            # Update options
            new_tenant = typer.prompt(f"New tenant (current: {user_info.get('tenant', 'N/A')}, press Enter to keep)", default="")
            new_role = typer.prompt(f"New role (current: {user_info.get('role', 'user')}, press Enter to keep)", default="")
            change_password = typer.confirm("Change password?")
            
            update_data = {}
            if new_tenant:
                update_data["tenant"] = new_tenant
            if new_role:
                update_data["role"] = new_role
            if change_password:
                new_password = typer.prompt("Enter new password", hide_input=True)
                update_data["password"] = new_password
            
            if update_data:
                if update_user(username, update_data):
                    console.print(f"[green]✅ User '{username}' updated successfully[/green]")
                else:
                    console.print(f"[red]❌ Failed to update user '{username}'[/red]")
            else:
                console.print("[yellow]No changes made[/yellow]")
            
            typer.prompt("Press Enter to continue")
        
        elif choice == "3":
            # Delete user
            console.print("\n[bold red]Delete User[/bold red]")
            if not users_db:
                console.print("[yellow]No users found[/yellow]")
                typer.prompt("Press Enter to continue")
                continue
            
            username = typer.prompt(f"Enter username to delete ({', '.join(users_db.keys())})")
            
            if username not in users_db:
                console.print(f"[red]❌ User '{username}' not found[/red]")
                typer.prompt("Press Enter to continue")
                continue
            
            if username == "admin":
                console.print("[red]❌ Cannot delete admin user[/red]")
                typer.prompt("Press Enter to continue")
                continue
            
            confirm = typer.confirm(f"Are you sure you want to delete user '{username}'?")
            if confirm:
                if delete_user(username):
                    console.print(f"[green]✅ User '{username}' deleted successfully[/green]")
                else:
                    console.print(f"[red]❌ Failed to delete user '{username}'[/red]")
            else:
                console.print("[yellow]Deletion cancelled[/yellow]")
            
            typer.prompt("Press Enter to continue")
        
        elif choice == "4":
            # Enable/Disable user
            console.print("\n[bold blue]Enable/Disable User[/bold blue]")
            if not users_db:
                console.print("[yellow]No users found[/yellow]")
                typer.prompt("Press Enter to continue")
                continue
            
            username = typer.prompt(f"Enter username ({', '.join(users_db.keys())})")
            
            if username not in users_db:
                console.print(f"[red]❌ User '{username}' not found[/red]")
                typer.prompt("Press Enter to continue")
                continue
            
            current_status = users_db[username].get("disabled", False)
            new_status = not current_status
            action = "disable" if new_status else "enable"
            
            confirm = typer.confirm(f"Are you sure you want to {action} user '{username}'?")
            if confirm:
                if update_user(username, {"disabled": new_status}):
                    console.print(f"[green]✅ User '{username}' {action}d successfully[/green]")
                else:
                    console.print(f"[red]❌ Failed to {action} user '{username}'[/red]")
            else:
                console.print(f"[yellow]{action.capitalize()} cancelled[/yellow]")
            
            typer.prompt("Press Enter to continue")
        
        elif choice == "5":
            # Back to main menu
            break
        
        else:
            console.print("[red]Invalid choice. Please try again.[/red]")
            typer.prompt("Press Enter to continue")

    def show_ingest():
        console.clear()
        console.print(Panel.fit("Content Ingestion", border_style="cyan"))
        
        console.print("[bold]Available Options:[/bold]")
        console.print("1. Ingest from Google Drive")
        console.print("2. Ingest from Sitemap")
        console.print("3. Ingest from Local Files")
        console.print("4. Back to Main Menu")
        
        choice = typer.prompt("Select an option")
        
        if choice == "1":
            tenant = typer.prompt("Enter tenant name")
            agent = typer.prompt("Enter agent name")
            folder_id = typer.prompt("Enter Google Drive folder ID")
            
            try:
                with console.status("Ingesting from Google Drive..."):
                    ingest(tenant, agent, drive=folder_id, console=console)
                console.print("[green]✅ Ingestion completed successfully[/green]")
            except Exception as e:
                console.print(f"[red]❌ Error: {str(e)}[/red]")
        
        elif choice == "2":
            tenant = typer.prompt("Enter tenant name")
            agent = typer.prompt("Enter agent name")
            sitemap_url = typer.prompt("Enter sitemap URL")
            
            try:
                with console.status("Ingesting from sitemap..."):
                    ingest(tenant, agent, sitemap=sitemap_url, console=console)
                console.print("[green]✅ Ingestion completed successfully[/green]")
            except Exception as e:
                console.print(f"[red]❌ Error: {str(e)}[/red]")
        
        elif choice == "3":
            tenant = typer.prompt("Enter tenant name")
            agent = typer.prompt("Enter agent name")
            file_path = typer.prompt("Enter file path")
            
            try:
                files = [Path(file_path)]
                with console.status("Ingesting local files..."):
                    ingest(tenant, agent, files=files, console=console)
                console.print("[green]✅ Ingestion completed successfully[/green]")
            except Exception as e:
                console.print(f"[red]❌ Error: {str(e)}[/red]")
        
        if choice != "4":
            typer.prompt("Press Enter to continue")
    
    def show_status():
        console.clear()
        console.print(Panel.fit("System Status", border_style="magenta"))
        
        # Check database
        db_status = "✅ Connected" if DB_PATH.exists() else "❌ Not Found"
        
        # Check API keys
        openai_key = "✅ Set" if os.getenv("OPENAI_API_KEY") else "❌ Missing"
        google_creds = "✅ Set" if os.getenv("GOOGLE_APPLICATION_CREDENTIALS") else "❌ Missing"
        
        # Count total interactions
        total_interactions = 0
        if DB_PATH.exists():
            with sqlite3.connect(DB_PATH) as con:
                cursor = con.execute("SELECT COUNT(*) FROM chat_logs")
                total_interactions = cursor.fetchone()[0]
        
        table = Table(title="System Status")
        table.add_column("Component")
        table.add_column("Status")
        table.add_row("Database", db_status)
        table.add_row("OpenAI API Key", openai_key)
        table.add_row("Google Credentials", google_creds)
        table.add_row("Total Interactions", str(total_interactions))
        
        console.print(table)
        typer.prompt("\nPress Enter to continue")
    
    # Main dashboard loop
    while True:
        try:
            choice = show_main_menu()
            
            if choice == "1":
                show_analytics()
            elif choice == "2":
                show_tenants()
            elif choice == "3":
                show_users()
            elif choice == "4":
                show_ingest()
            elif choice == "5":
                show_status()
            elif choice == "6":
                console.print("[green]Goodbye![/green]")
                break
            else:
                console.print("[red]Invalid choice. Please try again.[/red]")
                typer.prompt("Press Enter to continue")
                
        except KeyboardInterrupt:
            console.print("\n[yellow]Exiting...[/yellow]")
            break
        except Exception as e:
            console.print(f"[red]Error: {str(e)}[/red]")
            typer.prompt("Press Enter to continue")

@cli_app.command("serve")
def serve(
    host: str = typer.Option("0.0.0.0", help="Host to bind"),
    port: int = typer.Option(8000, help="Port to bind"),
    reload: bool = typer.Option(False, help="Enable auto-reload")
):
    """Start the web server"""
    import uvicorn
    
    console.print(f"🚀 Starting server on {host}:{port}")
    console.print(f"📊 Dashboard: http://{host}:{port}/docs")
    console.print(f"🤖 Widget: http://{host}:{port}/widget.js")
    
    uvicorn.run(
        "rag_chatbot:app",
        host=host,
        port=port,
        reload=reload
    )

@cli_app.command("ingest")
def cli_ingest(
    tenant: str = typer.Argument(help="Tenant name"),
    agent: str = typer.Argument(help="Agent name"),
    sitemap: Optional[str] = typer.Option(None, "--sitemap", "-s", help="Sitemap URL"),
    drive: Optional[str] = typer.Option(None, "--drive", "-d", help="Google Drive folder ID"),
    files: Optional[List[str]] = typer.Option(None, "--file", "-f", help="Local file paths")
):
    """Ingest content into vector store"""
    
    file_paths = [Path(f) for f in files] if files else None
    
    try:
        ingest(tenant, agent, sitemap=sitemap, drive=drive, files=file_paths, console=console)
        console.print("[green]✅ Ingestion completed successfully[/green]")
    except Exception as e:
        console.print(f"[red]❌ Error: {str(e)}[/red]")
        raise typer.Exit(1)

@cli_app.command("create-user")
def create_user_cli(
    username: str = typer.Argument(help="Username"),
    password: str = typer.Argument(help="Password"),
    tenant: str = typer.Option("*", help="Tenant (use '*' for all tenants)"),
    role: str = typer.Option("user", help="User role (user/admin)")
):
    """Create a new user"""
    
    user_data = UserCreate(
        username=username,
        password=password,
        tenant=tenant,
        role=role
    )
    
    if create_user(user_data):
        console.print(f"[green]✅ User '{username}' created successfully[/green]")
    else:
        console.print(f"[red]❌ User '{username}' already exists[/red]")
        raise typer.Exit(1)

if __name__ == "__main__":
    cli_app()

#added API endpoint for company specific agent listing 5/31/2025

@app.get("/my-agents")
async def get_my_agents(current_user: User = Depends(get_current_active_user)):
    """Get all agents available to the current user's tenant"""
    
    user_tenant = current_user.tenant
    
    # If user has access to all tenants, return all
    if user_tenant == "*":
        tenants = []
        for tenant_dir in BASE_CONFIG_DIR.iterdir():
            if tenant_dir.is_dir():
                agents = []
                for config_file in tenant_dir.iterdir():
                    if config_file.is_file() and config_file.suffix == ".json":
                        agent_name = config_file.stem
                        config = load_config(tenant_dir.name, agent_name)
                        agents.append({
                            "agent": agent_name,
                            "bot_name": config.get("bot_name", f"{tenant_dir.name}-{agent_name}"),
                            "primary_color": config.get("primary_color", "#1E88E5"),
                            "avatar_url": config.get("avatar_url", "")
                        })
                tenants.append({"tenant": tenant_dir.name, "agents": agents})
        return tenants
    
    # Return only agents for the user's specific tenant
    tenant_dir = BASE_CONFIG_DIR / user_tenant
    if not tenant_dir.exists():
        return {"tenant": user_tenant, "agents": []}
    
    agents = []
    for config_file in tenant_dir.iterdir():
        if config_file.is_file() and config_file.suffix == ".json":
            agent_name = config_file.stem
            config = load_config(user_tenant, agent_name)
            
            # Check if vector store exists for this agent
            vector_store_exists = store_path(user_tenant, agent_name).exists()
            
            agents.append({
                "agent": agent_name,
                "bot_name": config.get("bot_name", f"{user_tenant}-{agent_name}"),
                "primary_color": config.get("primary_color", "#1E88E5"),
                "secondary_color": config.get("secondary_color", "#FFFFFF"),
                "avatar_url": config.get("avatar_url", ""),
                "mode": config.get("mode", "inline"),
                "auto_open": config.get("auto_open", False),
                "vector_store_ready": vector_store_exists
            })
    
    return {
        "tenant": user_tenant,
        "agents": agents,
        "total_agents": len(agents)
    }

@app.get("/agent-status/{tenant}/{agent}")
async def get_agent_status(
    tenant: str,
    agent: str,
    current_user: User = Depends(get_current_active_user)
):
    """Get detailed status of a specific agent"""
    
    # Check if user has access to this tenant
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant"
        )
    
    # Check if agent exists
    config_file = cfg_path(tenant, agent)
    if not config_file.exists():
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Agent not found"
        )
    
    config = load_config(tenant, agent)
    vector_store_exists = store_path(tenant, agent).exists()
    
    # Get usage statistics from database
    with sqlite3.connect(DB_PATH) as con:
        cursor = con.execute(
            "SELECT COUNT(*) as total_chats, COUNT(DISTINCT session_id) as unique_sessions FROM chat_logs WHERE tenant = ? AND agent = ?",
            (tenant, agent)
        )
        total_chats, unique_sessions = cursor.fetchone()
        
        # Get recent activity (last 7 days)
        week_ago = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
        cursor = con.execute(
            "SELECT COUNT(*) FROM chat_logs WHERE tenant = ? AND agent = ? AND ts > ?",
            (tenant, agent, week_ago)
        )
        recent_chats = cursor.fetchone()[0]
    
    return {
        "tenant": tenant,
        "agent": agent,
        "config": config,
        "vector_store_ready": vector_store_exists,
        "statistics": {
            "total_chats": total_chats,
            "unique_sessions": unique_sessions,
            "recent_chats_7d": recent_chats
        }
    }

# Add this health endpoint

from datetime import datetime, timezone

@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "version": "7.0",
        "timestamp": datetime.now(timezone.utc).isoformat()
    }

# Add these enhanced endpoints to your rag_chatbot.py file
# REQUESTED A WEB INTERFACE

# Enhanced configuration model with widget parameters
class EnhancedConfigUpdateRequest(BaseModel):
    bot_name: Optional[str] = None
    system_prompt: Optional[str] = None
    primary_color: Optional[str] = None
    secondary_color: Optional[str] = None
    avatar_url: Optional[str] = None
    mode: Optional[str] = None
    auto_open: Optional[bool] = None
    llm_provider: Optional[str] = None
    llm_model: Optional[str] = None
    temperature: Optional[float] = None
    allowed_domains: Optional[List[str]] = None
    # New widget feature parameters
    enable_voice: Optional[bool] = None
    enable_files: Optional[bool] = None
    enable_tts: Optional[bool] = None
    enable_dark_mode: Optional[bool] = None
    widget_position: Optional[str] = None  # "bottom-right", "bottom-left", etc.
    widget_size: Optional[str] = None      # "small", "medium", "large"
    welcome_message: Optional[str] = None
    placeholder_text: Optional[str] = None

# Enhanced load_config function
def load_config(t: str, a: str) -> Dict[str, object]:
    p = cfg_path(t, a)
    if p.exists():
        return json.loads(p.read_text())
    p.parent.mkdir(parents=True, exist_ok=True)
    cfg = {
        "bot_name": f"{t}-{a}-Bot", 
        "system_prompt": "You are a helpful assistant.",
        "primary_color": "#1E88E5", 
        "secondary_color": "#FFFFFF", 
        "avatar_url": "",
        "mode": "inline", 
        "auto_open": False,
        "llm_provider": "openai",
        "llm_model": "gpt-4o-mini",
        "temperature": 0.3,
        "allowed_domains": ["*"],
        # Enhanced widget parameters
        "enable_voice": True,
        "enable_files": True,
        "enable_tts": False,
        "enable_dark_mode": True,
        "widget_position": "bottom-right",
        "widget_size": "medium",
        "welcome_message": "Hello! How can I help you today?",
        "placeholder_text": "Type your message..."
    }
    p.write_text(json.dumps(cfg, indent=2))
    return cfg

# Serve the web interface
@app.get("/admin", response_class=HTMLResponse)
async def get_admin_interface():
    """Serve the web administration interface"""
    # In production, you'd serve this from a static file
    # For now, we'll return a redirect to the HTML file
    return HTMLResponse("""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Redirecting...</title>
        <script>
            window.location.href = '/admin.html';
        </script>
    </head>
    <body>
        <p>Redirecting to admin interface...</p>
    </body>
    </html>
    """)

# Enhanced config endpoint with full widget parameters
@app.put("/config")
async def update_config_enhanced(
    config: EnhancedConfigUpdateRequest,
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT),
    current_user: User = Depends(get_admin_user)
):
    cfg = load_config(tenant, agent)
    
    # Update only provided fields
    for field, value in config.dict(exclude_unset=True).items():
        if value is not None:
            cfg[field] = value
    
    save_config(tenant, agent, cfg)
    return {"message": "Configuration updated successfully", "config": cfg}

# Get full config including widget parameters
@app.get("/config/full")
async def get_full_config(
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT),
    current_user: User = Depends(get_current_active_user)
):
    """Get full configuration including all widget parameters for admin interface"""
    
    # Check if user has access to this tenant
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant"
        )
    
    cfg = load_config(tenant, agent)
    return cfg

# Enhanced widget JavaScript with configurable parameters
def generate_enhanced_widget_js(tenant: str, agent: str) -> str:
    """Generate widget JavaScript with configuration-based parameters"""
    
    cfg = load_config(tenant, agent)
    
    return f"""(function(){{
const p=new URLSearchParams(location.search);
const tenant=p.get('tenant')||'{tenant}';
const agent=p.get('agent')||'{agent}';
const sid=sessionStorage.getItem('cq_sid')||(()=>{{const r=Math.random().toString(36).slice(2);sessionStorage.setItem('cq_sid',r);return r}})();
const msgs=[];
function $(id){{return document.getElementById(id);}}

// Configuration from server
const config = {json.dumps(cfg)};

// Create widget container
const container = document.createElement('div');
container.id = 'cq-widget-container';
document.body.appendChild(container);

// Features based on configuration
const features = {{
  typing: true,
  fileAttachments: config.enable_files || false,
  voiceInput: config.enable_voice || false,
  darkMode: config.enable_dark_mode || false,
  responseSpeech: config.enable_tts || false
}};

// Authentication mechanism
let authToken = localStorage.getItem('cq_auth_token');
const headers = {{
  'Content-Type': 'application/json',
  'X-Session-Id': sid
}};
if (authToken) headers['Authorization'] = `Bearer ${{authToken}}`;

function initWidget() {{
  // Dynamic positioning based on config
  const position = config.widget_position || 'bottom-right';
  const [vPos, hPos] = position.split('-');
  
  // Size configuration
  const size = config.widget_size || 'medium';
  const sizes = {{
    small: {{ width: '300px', height: '400px' }},
    medium: {{ width: '350px', height: '500px' }},
    large: {{ width: '400px', height: '600px' }}
  }};
  
  const widgetSize = sizes[size] || sizes.medium;
  
  // Create chat interface with dynamic configuration
  container.innerHTML = `
    <div id="cq-chat-widget" style="--primary:${{config.primary_color}};--secondary:${{config.secondary_color}};width:${{widgetSize.width}};height:${{widgetSize.height}}">
      <div id="cq-header">
        <img src="${{config.avatar_url || '/default-avatar.png'}}" alt="${{config.bot_name}}" onerror="this.src='data:image/svg+xml,<svg xmlns=%22http://www.w3.org/2000/svg%22 width=%2232%22 height=%2232%22><circle cx=%2216%22 cy=%2216%22 r=%2215%22 fill=%22%23ddd%22/><text x=%2216%22 y=%2220%22 text-anchor=%22middle%22 font-size=%2216%22>🤖</text></svg>'">
        <h3>${{config.bot_name}}</h3>
        <div class="cq-controls">
          ${{features.darkMode ? '<button id="cq-dark-toggle" title="Toggle dark mode">🌓</button>' : ''}}
          <button id="cq-settings">⚙️</button>
          <button id="cq-close">×</button>
        </div>
      </div>
      <div id="cq-messages">
        ${{config.welcome_message ? `<div class="cq-welcome-message">${{config.welcome_message}}</div>` : ''}}
      </div>
      <div id="cq-typing" style="display:none">
        <div class="cq-dot"></div>
        <div class="cq-dot"></div>
        <div class="cq-dot"></div>
      </div>
      <div id="cq-footer">
        <div id="cq-attachments"></div>
        <textarea id="cq-input" placeholder="${{config.placeholder_text || 'Type your message...'}}"></textarea>
        <div class="cq-actions">
          ${{features.voiceInput ? '<button id="cq-mic" title="Voice input">🎤</button>' : ''}}
          ${{features.fileAttachments ? '<button id="cq-attach" title="Attach file">📎</button>' : ''}}
          <button id="cq-send">➤</button>
        </div>
        ${{features.fileAttachments ? '<input type="file" id="cq-file-input" style="display:none" multiple accept=".pdf,.txt,.md,.docx">' : ''}}
      </div>
      <div id="cq-settings-panel" style="display:none">
        <h4>Settings</h4>
        ${{features.darkMode ? '<label><input type="checkbox" id="cq-dark-mode"> Dark mode</label>' : ''}}
        ${{features.responseSpeech ? '<label><input type="checkbox" id="cq-speech"> Read responses aloud</label>' : ''}}
        <h4>Sources</h4>
        <div id="cq-sources"></div>
      </div>
    </div>
  `;
  
  // Set up event handlers
  $('cq-send').addEventListener('click', sendMessage);
  $('cq-input').addEventListener('keypress', e => e.key === 'Enter' && !e.shiftKey && (e.preventDefault(), sendMessage()));
  $('cq-settings').addEventListener('click', toggleSettings);
  $('cq-close').addEventListener('click', minimizeWidget);
  
  // Conditional event handlers based on features
  if (features.voiceInput && $('cq-mic')) {{
    $('cq-mic').addEventListener('click', startVoiceInput);
  }}
  
  if (features.fileAttachments && $('cq-attach')) {{
    $('cq-attach').addEventListener('click', () => $('cq-file-input').click());
    $('cq-file-input').addEventListener('change', handleFileAttachment);
  }}
  
  if (features.darkMode && $('cq-dark-mode')) {{
    $('cq-dark-mode').addEventListener('change', e => setDarkMode(e.target.checked));
  }}
  
  if (features.darkMode && $('cq-dark-toggle')) {{
    $('cq-dark-toggle').addEventListener('click', toggleDarkMode);
  }}
  
  if (features.responseSpeech && $('cq-speech')) {{
    $('cq-speech').addEventListener('change', e => features.responseSpeech = e.target.checked);
  }}
  
  // Apply positioning
  container.style.position = 'fixed';
  container.style.zIndex = '10000';
  
  if (hPos === 'right') container.style.right = '20px';
  else container.style.left = '20px';
  
  if (vPos === 'bottom') container.style.bottom = '20px';
  else container.style.top = '20px';
  
  // Auto-open if configured
  if (config.auto_open) {{
    showWidget();
  }}
}}

function sendMessage() {{
  const input = $('cq-input');
  const text = input.value.trim();
  if (!text) return;
  
  addMessage('user', text);
  input.value = '';
  
  // Show typing indicator
  if (features.typing) $('cq-typing').style.display = 'flex';
  
  // Send to backend
  fetch(`/chat?tenant=${{tenant}}&agent=${{agent}}`, {{
    method: 'POST',
    headers: headers,
    body: JSON.stringify({{messages: msgs}})
  }})
  .then(r => r.json())
  .then(data => {{
    if (features.typing) $('cq-typing').style.display = 'none';
    addMessage('assistant', data.reply);
    updateSources(data.sources);
    
    // Text-to-speech if enabled
    if (features.responseSpeech) {{
      const speech = new SpeechSynthesisUtterance(data.reply);
      window.speechSynthesis.speak(speech);
    }}
  }})
  .catch(error => {{
    if (features.typing) $('cq-typing').style.display = 'none';
    addMessage('assistant', 'Sorry, I encountered an error. Please try again.');
    console.error('Chat error:', error);
  }});
}}

function addMessage(role, content) {{
  msgs.push({{role, content}});
  
  const msgDiv = document.createElement('div');
  msgDiv.className = `cq-message cq-${{role}}`;
  msgDiv.innerHTML = `
    <div class="cq-avatar">${{role === 'user' ? '👤' : '🤖'}}</div>
    <div class="cq-bubble">${{formatMessage(content)}}</div>
  `;
  
  $('cq-messages').appendChild(msgDiv);
  $('cq-messages').scrollTop = $('cq-messages').scrollHeight;
}}

function formatMessage(text) {{
  return text
    .replace(/\\*\\*(.*?)\\*\\*/g, '<strong>$1</strong>')
    .replace(/\\*(.*?)\\*/g, '<em>$1</em>')
    .replace(/```(.*?)```/gs, '<pre><code>$1</code></pre>')
    .replace(/`(.*?)`/g, '<code>$1</code>')
    .replace(/\\n/g, '<br>');
}}

function startVoiceInput() {{
  if (!('webkitSpeechRecognition' in window)) {{
    alert('Voice input not supported in your browser');
    return;
  }}
  
  const recognition = new webkitSpeechRecognition();
  recognition.continuous = false;
  recognition.interimResults = false;
  
  recognition.onstart = () => {{
    $('cq-mic').style.background = '#f44336';
  }};
  
  recognition.onend = () => {{
    $('cq-mic').style.background = '';
  }};
  
  recognition.onresult = function(event) {{
    const transcript = event.results[0][0].transcript;
    $('cq-input').value = transcript;
  }};
  
  recognition.start();
}}

function handleFileAttachment(event) {{
  const files = Array.from(event.target.files);
  if (files.length === 0) return;
  
  const attachmentsDiv = $('cq-attachments');
  attachmentsDiv.innerHTML = '';
  
  files.forEach(file => {{
    const fileDiv = document.createElement('div');
    fileDiv.className = 'cq-attachment';
    fileDiv.innerHTML = `
      <span>📎 ${{file.name}}</span>
      <button onclick="this.parentElement.remove()">×</button>
    `;
    attachmentsDiv.appendChild(fileDiv);
  }});
  
  // You would implement file upload to your backend here
  console.log('Files selected:', files);
}}

function updateSources(sources) {{
  const sourcesDiv = $('cq-sources');
  sourcesDiv.innerHTML = '';
  
  if (sources && sources.length > 0) {{
    const sourceList = document.createElement('ul');
    sources.forEach(src => {{
      const li = document.createElement('li');
      li.innerHTML = `<a href="${{src.source}}" target="_blank">${{src.source}}</a>`;
      sourceList.appendChild(li);
    }});
    sourcesDiv.appendChild(sourceList);
  }}
}}

function toggleSettings() {{
  const panel = $('cq-settings-panel');
  panel.style.display = panel.style.display === 'none' ? 'block' : 'none';
}}

function minimizeWidget() {{
  const widget = $('cq-chat-widget');
  widget.style.display = 'none';
  $('cq-launcher').style.display = 'block';
}}

function showWidget() {{
  const widget = $('cq-chat-widget');
  widget.style.display = 'flex';
  if ($('cq-launcher')) $('cq-launcher').style.display = 'none';
}}

function toggleWidget() {{
  const widget = $('cq-chat-widget');
  if (widget.style.display === 'none' || !widget.style.display) {{
    showWidget();
  }} else {{
    minimizeWidget();
  }}
}}

function setDarkMode(enabled) {{
  features.darkMode = enabled;
  const widget = $('cq-chat-widget');
  if (enabled) {{
    widget.classList.add('cq-dark');
  }} else {{
    widget.classList.remove('cq-dark');
  }}
  localStorage.setItem('cq_dark_mode', enabled);
}}

function toggleDarkMode() {{
  const current = localStorage.getItem('cq_dark_mode') === 'true';
  setDarkMode(!current);
  if ($('cq-dark-mode')) {{
    $('cq-dark-mode').checked = !current;
  }}
}}

// Enhanced CSS with more customization
const style = document.createElement('style');
style.textContent = `
  #cq-widget-container {{
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
  }}
  
  #cq-launcher {{
    width: 60px;
    height: 60px;
    border-radius: 50%;
    background: var(--primary, {cfg.get('primary_color', '#1E88E5')});
    color: white;
    border: none;
    font-size: 24px;
    cursor: pointer;
    box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    transition: all 0.3s ease;
  }}
  
  #cq-launcher:hover {{
    transform: scale(1.1);
    box-shadow: 0 6px 16px rgba(0,0,0,0.2);
  }}
  
  #cq-chat-widget {{
    border-radius: 12px;
    overflow: hidden;
    display: flex;
    flex-direction: column;
    box-shadow: 0 8px 32px rgba(0,0,0,0.12);
    background: white;
    transition: all 0.3s ease;
    border: 1px solid #e0e0e0;
  }}
  
  #cq-chat-widget.cq-dark {{
    background: #2d2d2d;
    border-color: #404040;
    color: white;
  }}
  
  #cq-header {{
    background: var(--primary, {cfg.get('primary_color', '#1E88E5')});
    color: white;
    padding: 16px;
    display: flex;
    align-items: center;
    gap: 12px;
  }}
  
  #cq-header img {{
    width: 32px;
    height: 32px;
    border-radius: 50%;
    object-fit: cover;
  }}
  
  #cq-header h3 {{
    margin: 0;
    flex: 1;
    font-size: 16px;
    font-weight: 600;
  }}
  
  .cq-controls {{
    display: flex;
    gap: 8px;
  }}
  
  .cq-controls button {{
    background: rgba(255,255,255,0.2);
    border: none;
    color: white;
    width: 32px;
    height: 32px;
    border-radius: 6px;
    cursor: pointer;
    transition: background 0.2s;
  }}
  
  .cq-controls button:hover {{
    background: rgba(255,255,255,0.3);
  }}
  
  #cq-messages {{
    flex: 1;
    padding: 16px;
    overflow-y: auto;
    display: flex;
    flex-direction: column;
    gap: 12px;
  }}
  
  .cq-welcome-message {{
    text-align: center;
    padding: 12px;
    background: #f0f7ff;
    border-radius: 8px;
    color: #1976d2;
    font-style: italic;
    margin-bottom: 12px;
  }}
  
  .cq-dark .cq-welcome-message {{
    background: #1a237e;
    color: #90caf9;
  }}
  
  .cq-message {{
    display: flex;
    gap: 8px;
    align-items: flex-start;
  }}
  
  .cq-message.cq-user {{
    flex-direction: row-reverse;
  }}
  
  .cq-avatar {{
    width: 32px;
    height: 32px;
    border-radius: 50%;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 14px;
    flex-shrink: 0;
  }}
  
  .cq-bubble {{
    max-width: 80%;
    padding: 12px 16px;
    border-radius: 18px;
    line-height: 1.4;
    word-wrap: break-word;
  }}
  
  .cq-user .cq-bubble {{
    background: var(--primary, {cfg.get('primary_color', '#1E88E5')});
    color: white;
  }}
  
  .cq-assistant .cq-bubble {{
    background: #f5f5f5;
    color: #333;
  }}
  
  .cq-dark .cq-assistant .cq-bubble {{
    background: #404040;
    color: white;
  }}
  
  #cq-typing {{
    padding: 16px;
    display: flex;
    justify-content: center;
    gap: 4px;
  }}
  
  .cq-dot {{
    width: 8px;
    height: 8px;
    border-radius: 50%;
    background: #ccc;
    animation: cq-bounce 1.4s infinite ease-in-out both;
  }}
  
  .cq-dot:nth-child(1) {{ animation-delay: -0.32s; }}
  .cq-dot:nth-child(2) {{ animation-delay: -0.16s; }}
  
  @keyframes cq-bounce {{
    0%, 80%, 100% {{ transform: scale(0); }}
    40% {{ transform: scale(1); }}
  }}
  
  #cq-footer {{
    padding: 16px;
    border-top: 1px solid #e0e0e0;
  }}
  
  .cq-dark #cq-footer {{
    border-top-color: #404040;
  }}
  
  #cq-attachments {{
    margin-bottom: 8px;
  }}
  
  .cq-attachment {{
    display: inline-flex;
    align-items: center;
    gap: 8px;
    background: #f0f0f0;
    padding: 4px 8px;
    border-radius: 16px;
    font-size: 12px;
    margin-right: 8px;
    margin-bottom: 4px;
  }}
  
  .cq-dark .cq-attachment {{
    background: #555;
    color: white;
  }}
  
  .cq-attachment button {{
    background: none;
    border: none;
    cursor: pointer;
    padding: 0;
    width: 16px;
    height: 16px;
    border-radius: 50%;
    display: flex;
    align-items: center;
    justify-content: center;
  }}
  
  #cq-input {{
    width: 100%;
    min-height: 20px;
    max-height: 100px;
    padding: 12px;
    border: 1px solid #e0e0e0;
    border-radius: 24px;
    resize: none;
    outline: none;
    font-family: inherit;
    font-size: 14px;
    line-height: 1.4;
  }}
  
  .cq-dark #cq-input {{
    background: #404040;
    border-color: #555;
    color: white;
  }}
  
  .cq-actions {{
    display: flex;
    justify-content: space-between;
    align-items: center;
    margin-top: 8px;
  }}
  
  .cq-actions button {{
    background: none;
    border: none;
    cursor: pointer;
    padding: 8px;
    border-radius: 50%;
    transition: background 0.2s;
  }}
  
  .cq-actions button:hover {{
    background: rgba(0,0,0,0.1);
  }}
  
  .cq-dark .cq-actions button:hover {{
    background: rgba(255,255,255,0.1);
  }}
  
  #cq-send {{
    background: var(--primary, {cfg.get('primary_color', '#1E88E5')}) !important;
    color: white !important;
  }}
  
  #cq-settings-panel {{
    position: absolute;
    top: 60px;
    right: 0;
    background: white;
    border: 1px solid #e0e0e0;
    border-radius: 8px;
    padding: 16px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.15);
    min-width: 200px;
    max-height: 300px;
    overflow-y: auto;
  }}
  
  .cq-dark #cq-settings-panel {{
    background: #2d2d2d;
    border-color: #404040;
    color: white;
  }}
  
  #cq-settings-panel h4 {{
    margin: 0 0 12px 0;
    font-size: 14px;
    font-weight: 600;
  }}
  
  #cq-settings-panel label {{
    display: flex;
    align-items: center;
    gap: 8px;
    margin-bottom: 8px;
    cursor: pointer;
    font-size: 14px;
  }}
  
  #cq-sources ul {{
    list-style: none;
    padding: 0;
    margin: 8px 0 0 0;
  }}
  
  #cq-sources li {{
    margin-bottom: 4px;
  }}
  
  #cq-sources a {{
    color: var(--primary, {cfg.get('primary_color', '#1E88E5')});
    text-decoration: none;
    font-size: 12px;
  }}
  
  #cq-sources a:hover {{
    text-decoration: underline;
  }}
  
  .cq-dark #cq-sources a {{
    color: #4fc3f7;
  }}
  
  /* Mobile responsive */
  @media (max-width: 480px) {{
    #cq-chat-widget {{
      width: calc(100vw - 40px) !important;
      height: calc(100vh - 40px) !important;
      max-height: 600px;
    }}
  }}
`;
document.head.appendChild(style);

// Initialize
if (document.readyState === 'loading') {{
  document.addEventListener('DOMContentLoaded', init);
}} else {{
  init();
}}

function init() {{
  // Create launcher button
  const launcher = document.createElement('button');
  launcher.id = 'cq-launcher';
  launcher.innerHTML = '💬';
  launcher.onclick = toggleWidget;
  launcher.style.display = 'block';
  container.appendChild(launcher);
  
  // Load saved preferences
  const savedDarkMode = localStorage.getItem('cq_dark_mode') === 'true';
  if (savedDarkMode && features.darkMode) {{
    setTimeout(() => setDarkMode(true), 100);
  }}
  
  // Initialize widget
  initWidget();
}}

}})();"""

# Enhanced widget endpoint with configuration
@app.get("/widget.js", response_class=PlainTextResponse)
async def get_enhanced_widget(
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(DEFAULT_AGENT)
):
    return generate_enhanced_widget_js(tenant, agent)

# Serve static admin interface
@app.get("/admin.html", response_class=HTMLResponse)
async def serve_admin_html():
    """Serve the admin interface HTML file"""
    # In production, you would serve this from a static file
    # For development, you can return the HTML directly or serve from file
    try:
        admin_html_path = Path("admin.html")
        if admin_html_path.exists():
            return HTMLResponse(admin_html_path.read_text())
        else:
            return HTMLResponse("""
            <!DOCTYPE html>
            <html>
            <head><title>Admin Interface Not Found</title></head>
            <body>
                <h1>Admin Interface</h1>
                <p>Please save the admin HTML interface as 'admin.html' in your project directory.</p>
                <p>You can access the API documentation at <a href="/docs">/docs</a></p>
            </body>
            </html>
            """)
    except Exception as e:
        return HTMLResponse(f"<h1>Error loading admin interface</h1><p>{str(e)}</p>")

# Enhanced analytics with widget-specific metrics
@app.get("/analytics/widget")
async def get_widget_analytics(
    tenant: str = Query(DEFAULT_TENANT),
    agent: str = Query(None),
    start_date: Optional[str] = None,
    end_date: Optional[str] = None,
    current_user: User = Depends(get_admin_user)
):
    """Get widget-specific analytics including feature usage"""
    
    if current_user.tenant != "*" and current_user.tenant != tenant:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="You don't have access to this tenant's analytics"
        )
    
    # Set default date range
    if not end_date:
        end_date = datetime.now(timezone.utc).isoformat()
    if not start_date:
        start_date = (datetime.now(timezone.utc) - timedelta(days=30)).isoformat()
    
    # Base query
    query = "SELECT * FROM chat_logs WHERE tenant = ? AND ts BETWEEN ? AND ?"
    params = [tenant, start_date, end_date]
    
    if agent:
        query += " AND agent = ?"
        params.append(agent)
    
    with sqlite3.connect(DB_PATH) as con:
        df = pd.read_sql_query(query, con, params=params)
    
    if df.empty:
        return {"message": "No data available for selected period"}
    
    # Widget-specific analytics
    analytics = {
        "total_interactions": len(df),
        "unique_sessions": df['session_id'].nunique(),
        "avg_session_length": df.groupby('session_id').size().mean(),
        "peak_hours": df.groupby(pd.to_datetime(df['ts']).dt.hour).size().to_dict(),
        "response_times": {
            "avg": df['latency'].mean(),
            "p95": df['latency'].quantile(0.95),
            "p99": df['latency'].quantile(0.99)
        },
        "user_satisfaction": {
            "avg_rating": df['user_feedback'].mean() if 'user_feedback' in df else None,
            "total_ratings": df['user_feedback'].count() if 'user_feedback' in df else 0
        }
    }
    
    return analytics