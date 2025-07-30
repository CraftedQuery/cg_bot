"""
utils/file_processors.py - File processing utilities
"""
from pathlib import Path
from typing import List, Dict, Tuple, Optional

from bs4 import BeautifulSoup

try:
    from pdf2image import convert_from_path
    import pytesseract
except Exception:
    convert_from_path = None
    pytesseract = None


def _find_heading(lines: List[str], start: int) -> Optional[str]:
    """Return the last Markdown heading before the start line."""
    for i in range(start, -1, -1):
        line = lines[i].strip()
        if line.startswith("#"):
            return line.lstrip("#").strip() or None
    return None


def _chunk_text_with_lines(
    text: str, *, page: Optional[int] = None, lines_per_chunk: int = 20
) -> Tuple[List[str], List[Dict]]:
    """Split text into line-based chunks and capture metadata."""
    lines = text.splitlines()
    chunks, metas = [], []

    for i in range(0, len(lines), lines_per_chunk):
        chunk_lines = lines[i : i + lines_per_chunk]
        chunk = "\n".join(chunk_lines).strip()
        if not chunk:
            continue
        meta = {
            "line": i + 1,
            "heading": _find_heading(lines, i),
        }
        if page is not None:
            meta["page"] = page
        chunks.append(chunk)
        metas.append(meta)

    return chunks, metas


def process_file(file_path: Path, filename: str) -> Tuple[List[str], List[Dict], bool]:
    """Process a file and extract text chunks and metadata.

    Returns a tuple of chunks, metadata and a flag indicating if OCR was used.
    """
    ext = file_path.suffix.lower()
    
    ocr_used = False
    try:
        if ext == ".pdf":
            chunks, metadatas, ocr_used = _process_pdf(file_path)
        else:
            if ext in {".txt", ".md"}:
                raw_text = _process_text(file_path)
            elif ext == ".docx":
                raw_text = _process_docx(file_path)
            elif ext == ".pptx":
                raw_text = _process_pptx(file_path)
            elif ext == ".xlsx":
                raw_text = _process_xlsx(file_path)
            elif ext == ".csv":
                raw_text = _process_csv(file_path)
            elif ext in {".html", ".htm"}:
                raw_text = _process_html(file_path)
            else:
                return [], [], False

            chunks, metadatas = _chunk_text_with_lines(raw_text)

        for m in metadatas:
            m["source"] = filename

        return chunks, metadatas, ocr_used
        
    except Exception as e:
        print(f"Error processing file {filename}: {str(e)}")
        return [], [], False


def _process_pdf(file_path: Path) -> Tuple[List[str], List[Dict], bool]:
    """Process PDF file into chunks with page metadata.

    If a page has no extractable text, OCR will be applied using pdf2image and
    pytesseract. Returns a flag indicating whether OCR was used.
    """
    try:
        import pypdf
    except ImportError:
        raise ImportError("pypdf package not installed")

    reader = pypdf.PdfReader(str(file_path))
    chunks, metas = [], []
    ocr_used = False

    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if not text or not text.strip():
            if convert_from_path and pytesseract:
                try:
                    images = convert_from_path(
                        str(file_path), first_page=idx, last_page=idx
                    )
                    if images:
                        text = pytesseract.image_to_string(images[0])
                        if text.strip():
                            ocr_used = True
                except Exception:
                    pass
        if not text or not text.strip():
            continue
        page_chunks, page_meta = _chunk_text_with_lines(text, page=idx)
        chunks.extend(page_chunks)
        metas.extend(page_meta)

    if not chunks and convert_from_path and pytesseract:
        try:
            images = convert_from_path(str(file_path))
            text = "\n".join(pytesseract.image_to_string(im) for im in images)
            if text.strip():
                ocr_used = True
                chunks, metas = _chunk_text_with_lines(text)
        except Exception:
            pass

    return chunks, metas, ocr_used


def _process_text(file_path: Path) -> str:
    """Process text file"""
    return file_path.read_text(errors='ignore')


def _process_docx(file_path: Path) -> str:
    """Process DOCX file"""
    try:
        import docx
    except ImportError:
        raise ImportError("python-docx package not installed")
    
    doc = docx.Document(str(file_path))
    return "\n".join(para.text for para in doc.paragraphs)


def _process_pptx(file_path: Path) -> str:
    """Process PPTX file"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx package not installed")
    
    prs = Presentation(str(file_path))
    text_parts = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_parts.append(shape.text)
    
    return "\n".join(text_parts)


def _process_xlsx(file_path: Path) -> str:
    """Process XLSX file"""
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl package not installed")
    
    wb = openpyxl.load_workbook(str(file_path))
    text_parts = []
    
    for sheet in wb.worksheets:
        for row in sheet.rows:
            row_text = []
            for cell in row:
                if cell.value:
                    row_text.append(str(cell.value))
            if row_text:
                text_parts.append("\t".join(row_text))
    
    return "\n".join(text_parts)


def _process_csv(file_path: Path) -> str:
    """Process CSV file"""
    import csv
    
    text_parts = []
    with open(str(file_path), 'r', errors='ignore') as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            text_parts.append(",".join(str(cell) for cell in row))
    
    return "\n".join(text_parts)


def _process_html(file_path: Path) -> str:
    """Process HTML file"""
    with open(str(file_path), 'r', errors='ignore') as htmlfile:
        soup = BeautifulSoup(htmlfile.read(), "html.parser")
        return soup.get_text("\n", strip=True)